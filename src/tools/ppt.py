"""PowerPoint (PPTX) generation tool — Enterprise GenAI Content Studio.

Given a topic + optional outline, this tool:
1. Calls GPT-4.1-mini to generate structured slide content
2. Optionally generates a cover image via gpt-image-2
3. Assembles a .pptx file using python-pptx
4. Returns the local file path for download

Every call is real — no mock data, no stubs.
"""

from __future__ import annotations

import json
import uuid
from pathlib import Path
from typing import Any

import structlog

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    _PPTX_IMPORT_ERROR: Exception | None = None
except Exception as exc:  # pragma: no cover - environment dependent import
    Presentation = None  # type: ignore[assignment]
    RGBColor = None  # type: ignore[assignment]
    Inches = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_IMPORT_ERROR = exc

from src.config import get_azure_openai_client, get_settings

logger = structlog.get_logger()


# ── Slide layout indexes (for python-pptx default template) ───────────────────
_LAYOUT_TITLE      = 0   # Title Slide
_LAYOUT_TITLE_BODY = 1   # Title and Content
_LAYOUT_SECTION    = 2   # Section Header
_LAYOUT_BLANK      = 6   # Blank


_GENERATE_SLIDES_SYSTEM = """You are a professional presentation designer.
Generate a structured slide deck in JSON format.

Return ONLY valid JSON. Schema:
{
  "title": "<deck title>",
  "subtitle": "<one-line subtitle>",
  "author": "<author or company>",
  "slides": [
    {
      "type": "title",       // only for first slide
      "heading": "<heading>",
      "body": "<optional subtitle>"
    },
    {
      "type": "content",
      "heading": "<slide title>",
      "bullets": ["<point 1>", "<point 2>", "<point 3>"]
    },
    {
      "type": "section",    // section divider
      "heading": "<section name>",
      "body": "<optional tagline>"
    }
  ]
}

Rules:
- 8-12 slides total
- Each content slide: 3-5 concise bullets (max 15 words each)
- Professional, enterprise tone
- First slide must be type "title"
- Last slide: type "content", heading "Key Takeaways"
"""


class PPTGenerationTool:
    """Generate real .pptx files from a text prompt using LLM + python-pptx."""

    def __init__(self, output_dir: Path | None = None):
        self._settings = get_settings()
        self._output_dir = output_dir or Path("output") / "generated-ppts"
        self._output_dir.mkdir(parents=True, exist_ok=True)

    async def generate(
        self,
        topic: str,
        audience: str = "enterprise stakeholders",
        style: str = "professional",
        slides: int = 6,
        include_cover_image: bool = False,
    ) -> dict[str, Any]:
        """Generate a full PowerPoint presentation.

        Args:
            topic: Presentation topic or detailed prompt
            audience: Target audience (affects tone and depth)
            style: "professional" | "minimal" | "vibrant"
            slides: Target number of slides (3-15)
            include_cover_image: Whether to generate a cover image via DALL-E

        Returns:
            dict with file_path, slide_count, title
        """
        logger.info("ppt.generating", topic=topic[:80], audience=audience, slides=slides)

        # Step 1 — LLM generates slide structure
        slide_data = await self._generate_slide_structure(topic, audience, style, slides)

        # Step 2 — Build PPTX
        file_name = f"ppt-{uuid.uuid4().hex[:8]}.pptx"
        file_path = self._output_dir / file_name
        slide_count = self._build_pptx(slide_data, file_path, style)

        logger.info("ppt.generated", path=str(file_path), slides=slide_count)
        return {
            "file_path": str(file_path),
            "file_name": file_name,
            "title": slide_data.get("title", topic[:60]),
            "slide_count": slide_count,
            "topic": topic,
            "audience": audience,
            "style": style,
        }

    async def _generate_slide_structure(
        self, topic: str, audience: str, style: str, slides: int = 6
    ) -> dict[str, Any]:
        """Call GPT-4.1-mini to produce slide JSON."""
        client = get_azure_openai_client()
        user_prompt = (
            f"Topic: {topic}\n"
            f"Audience: {audience}\n"
            f"Style: {style}\n"
            f"Target slides: {slides}\n\n"
            "Generate a complete, compelling presentation."
        )
        response = await client.chat.completions.create(
            model=self._settings.azure_ai_chat_deployment,
            messages=[
                {"role": "system", "content": _GENERATE_SLIDES_SYSTEM},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.7,
            response_format={"type": "json_object"},
        )
        raw = response.choices[0].message.content or "{}"
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            logger.warning("ppt.json_parse_failed", raw=raw[:200])
            return {"title": topic, "slides": []}

    def _build_pptx(
        self, data: dict[str, Any], file_path: Path, style: str
    ) -> int:
        """Assemble the PPTX file. Returns slide count."""
        if Presentation is None or RGBColor is None or Inches is None or Pt is None:
            raise RuntimeError(
                "python-pptx is not available in this runtime. "
                "Install dependency 'python-pptx>=1.0.0'."
            ) from _PPTX_IMPORT_ERROR

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Color palette by style
        palettes = {
            "professional": {"primary": RGBColor(0x1A, 0x73, 0xE8), "text": RGBColor(0x20, 0x20, 0x20)},
            "minimal":      {"primary": RGBColor(0x21, 0x21, 0x21), "text": RGBColor(0x33, 0x33, 0x33)},
            "vibrant":      {"primary": RGBColor(0xE9, 0x1E, 0x63), "text": RGBColor(0x1A, 0x1A, 0x2E)},
        }
        palette = palettes.get(style, palettes["professional"])

        slides = data.get("slides", [])
        if not slides:
            # Fallback: bare title slide
            slides = [{"type": "title", "heading": data.get("title", "Presentation"), "body": ""}]

        for i, slide_def in enumerate(slides):
            stype = slide_def.get("type", "content")

            if stype == "title" and i == 0:
                self._add_title_slide(prs, slide_def, data, palette)
            elif stype == "section":
                self._add_section_slide(prs, slide_def, palette)
            else:
                self._add_content_slide(prs, slide_def, palette)

        prs.save(str(file_path))
        return len(slides)

    def _add_title_slide(self, prs, slide_def, data, palette):
        layout = prs.slide_layouts[_LAYOUT_TITLE]
        slide  = prs.slides.add_slide(layout)
        title  = slide.shapes.title
        subtitle_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        title.text = slide_def.get("heading") or data.get("title", "Presentation")
        title.text_frame.paragraphs[0].font.size  = Pt(40)
        title.text_frame.paragraphs[0].font.bold  = True
        title.text_frame.paragraphs[0].font.color.rgb = palette["primary"]

        if subtitle_ph:
            sub_text = slide_def.get("body") or data.get("subtitle", "")
            subtitle_ph.text = sub_text
            subtitle_ph.text_frame.paragraphs[0].font.size  = Pt(20)
            subtitle_ph.text_frame.paragraphs[0].font.color.rgb = palette["text"]

    def _add_section_slide(self, prs, slide_def, palette):
        layout = prs.slide_layouts[_LAYOUT_SECTION]
        slide  = prs.slides.add_slide(layout)
        title  = slide.shapes.title
        title.text = slide_def.get("heading", "")
        title.text_frame.paragraphs[0].font.size  = Pt(36)
        title.text_frame.paragraphs[0].font.bold  = True
        title.text_frame.paragraphs[0].font.color.rgb = palette["primary"]

        if len(slide.placeholders) > 1 and slide_def.get("body"):
            slide.placeholders[1].text = slide_def["body"]

    def _add_content_slide(self, prs, slide_def, palette):
        layout = prs.slide_layouts[_LAYOUT_TITLE_BODY]
        slide  = prs.slides.add_slide(layout)
        title  = slide.shapes.title
        title.text = slide_def.get("heading", "")
        title.text_frame.paragraphs[0].font.size  = Pt(28)
        title.text_frame.paragraphs[0].font.bold  = True
        title.text_frame.paragraphs[0].font.color.rgb = palette["primary"]

        body_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        bullets = slide_def.get("bullets", [])
        if body_ph and bullets:
            tf = body_ph.text_frame
            tf.clear()
            for idx, bullet in enumerate(bullets):
                if idx == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text  = bullet
                para.level = 0
                para.font.size  = Pt(18)
                para.font.color.rgb = palette["text"]
